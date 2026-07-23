#!/usr/bin/env python3
"""Original PPTX helper for structural inspection and common QA warnings."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


def shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()


def inspect(path: Path) -> None:
    deck = Presentation(path)
    slides = []
    for index, slide in enumerate(deck.slides, start=1):
        texts = [shape_text(shape) for shape in slide.shapes]
        nonempty_texts = [text for text in texts if text]
        slides.append(
            {
                "slide": index,
                "shapes": len(slide.shapes),
                "text_shapes": len(nonempty_texts),
                "images": sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE),
                "title": (slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else None),
                "text_preview": " | ".join(nonempty_texts)[:300],
            }
        )
    report = {
        "file": str(path.resolve()),
        "slides": len(deck.slides),
        "width_inches": round(deck.slide_width / 914400, 3),
        "height_inches": round(deck.slide_height / 914400, 3),
        "slide_details": slides,
    }
    print(json.dumps(report, ensure_ascii=False, indent=2))


def validate(path: Path, minimum_font_pt: float) -> None:
    deck = Presentation(path)
    warnings = []
    for slide_index, slide in enumerate(deck.slides, start=1):
        meaningful = 0
        for shape_index, shape in enumerate(slide.shapes, start=1):
            text = shape_text(shape)
            if text or shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                meaningful += 1
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > deck.slide_width or shape.top + shape.height > deck.slide_height:
                warnings.append({"slide": slide_index, "shape": shape_index, "kind": "outside_slide"})
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size < Pt(minimum_font_pt):
                            warnings.append(
                                {
                                    "slide": slide_index,
                                    "shape": shape_index,
                                    "kind": "small_font",
                                    "font_pt": round(run.font.size.pt, 1),
                                    "text": run.text[:80],
                                }
                            )
        if meaningful == 0:
            warnings.append({"slide": slide_index, "kind": "blank_slide"})
        title = slide.shapes.title
        if title and title.has_text_frame and "\n" in title.text.strip():
            warnings.append({"slide": slide_index, "kind": "multiline_title", "text": title.text.strip()[:160]})
    print(json.dumps({"file": str(path.resolve()), "warnings": warnings, "count": len(warnings)}, ensure_ascii=False, indent=2))


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser()
    sub = parser.add_subparsers(dest="command", required=True)
    inspect_cmd = sub.add_parser("inspect")
    inspect_cmd.add_argument("input", type=Path)
    validate_cmd = sub.add_parser("validate")
    validate_cmd.add_argument("input", type=Path)
    validate_cmd.add_argument("--minimum-font-pt", type=float, default=14.0)
    return parser


def main() -> None:
    args = build_parser().parse_args()
    if args.command == "inspect":
        inspect(args.input)
    elif args.command == "validate":
        validate(args.input, args.minimum_font_pt)


if __name__ == "__main__":
    main()
